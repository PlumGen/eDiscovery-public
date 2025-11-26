import os
import io
import zipfile
import tarfile
import tempfile
import py7zr
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import docx
import openpyxl
import csv
import extract_msg
import email
from striprtf.striprtf import rtf_to_text


class DocumentExtractor:
    def __init__(self):
        self.supported_archives = [".zip", ".tar", ".gz", ".tgz", ".7z"]
        self.results = []

    def extract_text(self, blob: bytes, filename: str) -> dict:
        """
        Extract text from a blob (bytes) given its filename.
        Returns dict {filename: extracted_text}.
        Handles recursion into nested archives.
        """
        ext = os.path.splitext(filename)[1].lower()

        if ext in self.supported_archives:
            return self._extract_archive(blob, filename)

        try:
            text = self._dispatch_extractor(blob, filename, ext)
            self.results.append({'filename':filename,
                    'Body': text})
        except Exception as e:
            self.results.append({'filename':filename,
                    'Body': f"[ERROR: {str(e)}]"})

    # ---------------- INTERNAL HELPERS ---------------- #

    def _dispatch_extractor(self, blob, filename, ext):
        if ext == ".pdf":
            return self._extract_pdf(blob)

        elif ext == ".docx":
            return self._extract_docx(blob)

        elif ext in [".doc"]:  # legacy Word (binary)
            return self._extract_doc(blob, filename)

        elif ext in [".xlsx", ".csv", ".xls"]:
            return self._extract_spreadsheet(blob, filename)

        elif ext == ".rtf":
            return self._extract_rtf(blob)

        elif ext == ".pptx":
            return self._extract_pptx(blob, filename)

        elif ext == ".eml":
            return self._extract_eml(blob)

        elif ext == ".msg":
            return self._extract_msg(blob, filename)

        elif ext in [".png", ".jpg", ".jpeg", ".tiff"]:
            return self._extract_image(blob)

        else:
            return f"[Unsupported file type: {ext}]"

    # ---------------- ARCHIVE HANDLING ---------------- #

    def _extract_archive(self, blob, filename):

        with tempfile.TemporaryDirectory() as tmpdir:
            # safe filename for writing
            safe_name = os.path.basename(filename)
            path = os.path.join(tmpdir, safe_name)

            os.makedirs(os.path.dirname(path), exist_ok=True)
            
            print('path')
            print(path)
            with open(path, "wb") as f:
                f.write(blob)

            ext = os.path.splitext(safe_name)[1].lower()

            if ext == ".zip":
                with zipfile.ZipFile(path, "r") as zf:
                    for inner in zf.namelist():
                        if zf.getinfo(inner).is_dir():
                            continue
                        inner_bytes = zf.read(inner)
                        # prepend parent archive name
                        nested_name = f"{filename}/{inner}"
                        self.extract_text(inner_bytes, nested_name)

            elif ext in [".tar", ".gz", ".tgz"]:
                with tarfile.open(path, "r:*") as tf:
                    for member in tf.getmembers():
                        if not member.isfile():
                            continue
                        inner_bytes = tf.extractfile(member).read()
                        nested_name = f"{filename}/{member.name}"
                        self.extract_text(inner_bytes, nested_name)

            elif ext == ".7z":
                with py7zr.SevenZipFile(path, mode="r") as archive:
                    with tempfile.TemporaryDirectory() as extract_dir:
                        archive.extractall(path=extract_dir)
                        for root, _, files in os.walk(extract_dir):
                            for inner in files:
                                full_path = os.path.join(root, inner)
                                with open(full_path, "rb") as f:
                                    filedata = f.read()
                                rel_path = os.path.relpath(full_path, extract_dir)
                                nested_name = f"{filename}/{rel_path.replace(os.sep, '/')}"
                                self.extract_text(filedata, nested_name)

        



    # ---------------- FORMAT EXTRACTORS ---------------- #

    def _extract_pdf(self, blob):
        text = ""
        doc = fitz.open(stream=blob, filetype="pdf")
        for page in doc:
            page_text = page.get_text()
            if not page_text.strip():  # scanned page → OCR
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                page_text = pytesseract.image_to_string(img)
            text += page_text + "\n"
        return text

    def _extract_docx(self, blob):
        doc = docx.Document(io.BytesIO(blob))
        return "\n".join([p.text for p in doc.paragraphs])

    def _extract_doc(self, blob, filename):
        # For .doc (binary Word) – use antiword if available, else fallback
        with tempfile.NamedTemporaryFile(delete=False, suffix=filename) as tmp:
            tmp.write(blob)
            tmp.flush()
            try:
                import subprocess
                output = subprocess.check_output(["antiword", tmp.name])
                return output.decode("utf-8", errors="ignore")
            except Exception:
                return "[.doc parsing not supported without antiword]"

    def _extract_spreadsheet(self, blob, filename):
        text = ""
        if filename.endswith(".csv"):
            decoded = io.StringIO(blob.decode("utf-8", errors="ignore"))
            reader = csv.reader(decoded)
            for row in reader:
                text += " ".join(row) + "\n"
        else:  # .xlsx or .xls
            with tempfile.NamedTemporaryFile(delete=False, suffix=filename) as tmp:
                tmp.write(blob)
                tmp.flush()
                wb = openpyxl.load_workbook(tmp.name, data_only=True)
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        text += " ".join([str(c) if c is not None else "" for c in row]) + "\n"
        return text

    def _extract_rtf(self, blob):
        return rtf_to_text(blob.decode("utf-8", errors="ignore"))

    def _extract_pptx(self, blob, filename):
        from pptx import Presentation
        with tempfile.NamedTemporaryFile(delete=False, suffix=filename) as tmp:
            tmp.write(blob)
            tmp.flush()
            prs = Presentation(tmp.name)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def _extract_eml(self, blob):
        decoded = io.StringIO(blob.decode("utf-8", errors="ignore"))
        msg = email.message_from_file(decoded)
        body = []
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    body.append(part.get_payload(decode=True).decode(errors="ignore"))
        else:
            body.append(msg.get_payload(decode=True).decode(errors="ignore"))
        return "\n".join(body)

    def _extract_msg(self, blob, filename):
        with tempfile.NamedTemporaryFile(delete=False, suffix=filename) as tmp:
            tmp.write(blob)
            tmp.flush()
            msg = extract_msg.Message(tmp.name)
            return "\n".join([msg.sender, msg.subject, msg.body])

    def _extract_image(self, blob):
        img = Image.open(io.BytesIO(blob))
        return pytesseract.image_to_string(img)
